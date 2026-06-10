from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/zhangzhao/PycharmProjects/shuangtan2026")
TEMPLATE_BACKUP = ROOT / "work_assets" / "附件2原模板备份.pptx"
CURRENT = ROOT / "附件2天津大学国际工程师学院双碳主题企业实战项目课答辩模板.pptx"
DONE = ROOT / "附件2天津大学国际工程师学院双碳主题企业实战项目课答辩-已完成.pptx"
OLD_BUSY = ROOT / "work_assets" / "附件2上一版复杂背景备份.pptx"
CLEAN_BACKUP = ROOT / "work_assets" / "附件2上一版清爽模板备份.pptx"
MEDIA = ROOT / "work_assets" / "template_media" / "ppt" / "media"
VISUALS = ROOT / "work_assets" / "ppt_clean_images"

BLUE = RGBColor(0, 83, 137)
BLUE_DARK = RGBColor(9, 43, 89)
BLUE_LIGHT = RGBColor(220, 236, 248)
TEXT = RGBColor(20, 40, 70)
MUTED = RGBColor(90, 105, 125)
GREY = RGBColor(244, 247, 250)
GREEN = RGBColor(36, 139, 111)
ORANGE = RGBColor(228, 126, 45)


SLIDES = [
    ("基于本地大模型与 RAG 技术\n的网站智能客服系统", "“双碳”主题企业实战项目课答辩汇报", "2026 年 6 月 12 日"),
    ("OUTLINE", "汇报内容", ["PROJECT INTRODUCTION", "PROJECT PROCESS", "PROJECT RESULT"]),
    ("1 PROJECT INTRODUCTION", "项目背景：企业技术资料需要可解释问答", ["资料分散：NXP AIoT Cloud / Edge AI 资料覆盖平台、工具链、模型和应用范例", "提问方式自然化：老师或用户不会严格按照固定 FAQ 关键词提问", "普通大模型缺少项目内资料约束，容易给出无法追溯的泛化回答", "本地 RAG 把答案绑定到知识库片段，并展示标题、来源和相似度", "双碳 / 边缘 AI 场景强调本地化、低依赖和可部署的智能服务"]),
    ("1 PROJECT INTRODUCTION", "需求与目标", ["网页端提问、示例问题、Top-K 控制和检索来源展示", "健康检查、统计、索引重建、聊天问答等后端接口完整可用", "回答同时返回内容、来源、相似度、流程状态和耗时指标", "默认本地 Ollama，不依赖云端大模型 API，适合课堂离线演示", "通过 AI 创作平台辅助搭建 Agent / RAG 工作流，核心数据与索引落在本地", "模型或向量库异常时可降级，不让演示流程中断"]),
    ("2 PROJECT PROCESS", "总体架构", ["前端交互层：HTML / CSS / JavaScript 三栏页面，负责提问、状态和来源展示", "后端服务层：FastAPI 接口、Pydantic schema、流程状态和异常处理", "RAG 检索层：文档加载、chunk、embedding、ChromaDB、混合召回和来源聚合", "本地模型层：Ollama qwen2.5:7b 生成回答，nomic-embed-text 生成向量", "数据层：Markdown 知识库和本地向量索引均保存在项目目录内"]),
    ("2 PROJECT PROCESS", "本地知识库构建", ["采集 NXP AIoT Cloud、Cloud Lab、AI Hub、Model Zoo 等资料", "整理为 17 篇 Markdown 文档，包含标题、分类、来源、摘要和正文", "按 chunk 切分为 134 个片段，并保存 metadata 便于追溯来源", "向量化后写入 ChromaDB；不可用时自动回退 JSON + numpy 检索", "覆盖 LLM/VLM Edge Studio、Qwen/VSS、YOLOv8n、i.MX/MCU 等主题"]),
    ("2 PROJECT PROCESS", "RAG 流程与混合检索", ["接收问题 → 向量化 → 检索 → Prompt → 本地模型生成 → 返回答案", "向量检索负责语义相似度，适合长问题和概念解释类问题", "关键词 / 标题匹配补强专有名词、短中文问题和资料数量类问题", "按文档聚合来源，避免同一文档多个 chunk 重复刷屏", "无明确依据时拒答并清空无关来源，减少“看起来有引用”的误导"]),
    ("2 PROJECT PROCESS", "本地模型与降级机制", ["LLM：qwen2.5:7b，用于生成面向用户的客服回答", "Embedding：nomic-embed-text，用于问题和文档片段向量化", "Vector Store：优先 ChromaDB 持久化索引，失败时退回 SimpleVectorStore", "Embedding 失败时退回 TF-IDF，保证基础检索流程仍可演示", "模型生成失败时可启用 MOCK_LLM，前端仍能展示 RAG 流程和来源"]),
    ("2 PROJECT PROCESS", "前端演示页面", ["左侧：Ollama 状态、模型名、文档数、chunk 数、向量库类型和索引重建", "中间：聊天消息、示例问题、Top-K 控制、输入框和发送状态", "右侧：RAG 六步流程、检索来源、相似度、类别和内容片段", "长问题、示例问题和回答区域均做换行约束，避免投屏时超出屏幕", "页面不依赖 React / Vue / CDN，后端启动后即可访问"]),
    ("2 PROJECT PROCESS", "后端 API 与代码模块", ["GET /api/health：模型连接、知识库数量、chunk 数和向量库类型", "GET /api/stats：文档数、分类分布、chunk 数和知识库范围", "POST /api/rebuild-index：清理旧索引并重新构建本地向量库", "POST /api/chat：执行检索、Prompt 构造、本地模型回答和来源返回", "模块分层：config、schemas、ollama_client、rag、frontend 便于维护"]),
    ("3 PROJECT RESULT", "测试与验收结果", ["pytest：8 项单元测试通过，覆盖配置、分块、检索和接口基础行为", "smoke_test：覆盖首页、健康检查、统计、重建索引和 MOCK 聊天", "acceptance_check：覆盖目录、配置、文档、前端元素、降级和来源优先级", "真实接口批测：资料数、RAG 流程、NXP 技术问题、越界问题和追问表达", "验收重点：不是只看回答，还要看来源、流程、耗时和降级状态"]),
    ("3 PROJECT RESULT", "问题修复与迭代优化", ["短问法召回偏差：加入混合召回和标题覆盖评分", "资料数量问题死板：改为运行上下文模式，不写硬回答模板", "重复来源：按文档聚合 sources，减少同标题 chunk 重复出现", "索引重建异常：清理旧索引，修复 Chroma readonly 场景", "前端超屏：补齐长文本换行、响应式宽度和输入区约束"]),
    ("3 PROJECT RESULT", "项目成果与价值", ["完成 FastAPI + Ollama + RAG + 原生前端完整项目", "形成 17 篇知识库文档、134 个 chunks 和 Chroma 向量索引", "回答带来源、相似度和处理流程，区别于普通大模型套壳", "本地化部署减少云端依赖，适合资料问答、课程演示和边缘 AI 场景", "形成 README、API 文档、设计文档、测试脚本、报告和答辩材料", "后续可扩展上传资料、增量索引、多轮会话和用户反馈闭环"]),
    ("THANK YOU !", "总结与展望", ["本项目验证了本地 RAG 智能客服的技术可行性", "已经覆盖知识库构建、向量索引、问答生成、来源追溯和降级演示", "后续可加入知识库上传、增量索引、多轮会话和用户反馈", "可进一步接入真实边缘设备资料和端侧 AI 演示", "感谢各位老师指导，欢迎批评指正"]),
]

TECH_NOTES = {
    3: "技术落点：用本地知识库约束回答范围，答案必须能回到资料标题、来源和检索片段。",
    4: "验收口径：能问、能检索、能生成、能追溯、能降级、能重新构建索引。",
    5: "代码入口：app/main.py；核心模块：retriever.py / prompt_builder.py / vector_store.py / ollama_client.py。",
    6: "入库参数：17 docs / 134 chunks；CHUNK_SIZE=500，CHUNK_OVERLAP=80，TOP_K 默认 3。",
    7: "召回策略：向量相似度 + 关键词 / 标题加权；低相关结果不进入最终来源。",
    8: "本机配置策略：qwen2.5:7b 负责生成，nomic-embed-text 负责向量化；异常时保持可演示闭环。",
    9: "交互细节：示例问题、Top-K、状态面板和来源片段同步刷新，长文本强制换行避免超屏。",
    10: "接口响应包含 answer、sources、metrics、steps 等字段，前端可直接渲染流程状态。",
    11: "验证组合：单元测试 + smoke test + acceptance check + 真实问答批测。",
    12: "迭代原则：先保证演示可用，再优化召回质量，最后处理边界问题与前端展示。",
    13: "项目交付：代码仓库、知识库、向量索引、运行脚本、报告、PPT 和逐页讲解稿。",
}


def delete_all_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        el = shape._element
        el.getparent().remove(el)


def fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()


def add_text(slide, x, y, w, h, text, size=24, bold=False, color=TEXT, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_logo(slide, x=0.35, y=0.18, w=1.55):
    logo = MEDIA / "image12.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(x), Inches(y), width=Inches(w))


def add_anniversary(slide, x=11.55, y=0.12, w=0.85):
    img = MEDIA / "image13.jpeg"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(x), Inches(y), width=Inches(w))


def add_header(slide, idx: int, section: str, title: str):
    add_logo(slide, 0.32, 0.16, 1.45)
    add_anniversary(slide, 11.62, 0.12, 0.75)
    add_text(slide, Inches(0.35), Inches(0.93), Inches(5.8), Inches(0.35), section, 12, True, BLUE)
    add_text(slide, Inches(0.35), Inches(1.27), Inches(8.2), Inches(0.55), title, 24, True, BLUE_DARK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(1.84), Inches(11.9), Inches(0.03))
    fill(line, BLUE)
    add_text(slide, Inches(11.72), Inches(6.97), Inches(0.55), Inches(0.25), f"{idx:02d}", 8, False, MUTED, PP_ALIGN.RIGHT)


def add_bullet_list(slide, bullets, x=0.8, y=2.1, w=5.3, h=3.8, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(7 if len(bullets) <= 4 else 5)
    return box


def add_card(slide, x, y, w, h, title, body="", accent=BLUE):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.04), Inches(y + 0.04), Inches(w), Inches(h))
    fill(shadow, RGBColor(220, 228, 238), 35)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(rect, RGBColor(255, 255, 255))
    rect.line.color.rgb = RGBColor(222, 230, 238)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    fill(bar, accent)
    compact = h < 0.68
    title_text = f"{title}｜{body}" if compact and body else title
    add_text(
        slide,
        Inches(x + 0.22),
        Inches(y + (0.11 if compact else 0.16)),
        Inches(w - 0.35),
        Inches(0.22 if compact else 0.32),
        title_text,
        9.2 if compact else 14,
        True,
        BLUE_DARK,
    )
    if body and not compact:
        add_text(slide, Inches(x + 0.22), Inches(y + 0.55), Inches(w - 0.35), Inches(h - 0.65), body, 10.5, False, MUTED)


def add_picture_card(slide, img_path: Path, x, y, w, h):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.05), Inches(y + 0.05), Inches(w), Inches(h))
    fill(shadow, RGBColor(220, 228, 238), 40)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(frame, RGBColor(255, 255, 255))
    frame.line.color.rgb = RGBColor(222, 230, 238)
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(x + 0.08), Inches(y + 0.08), width=Inches(w - 0.16), height=Inches(h - 0.16))
    else:
        add_text(slide, Inches(x + 0.25), Inches(y + 0.3), Inches(w - 0.5), Inches(0.4), img_path.name, 10, False, MUTED)


def add_label_strip(slide, labels, x, y, chip_w, chip_h=0.34):
    for i, label in enumerate(labels):
        color = [BLUE, GREEN, ORANGE, RGBColor(105, 91, 180), BLUE_DARK, GREEN][i % 6]
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + i * (chip_w + 0.12)),
            Inches(y),
            Inches(chip_w),
            Inches(chip_h),
        )
        fill(chip, RGBColor(246, 249, 252))
        chip.line.color.rgb = color
        add_text(
            slide,
            Inches(x + i * (chip_w + 0.12) + 0.04),
            Inches(y + 0.06),
            Inches(chip_w - 0.08),
            Inches(0.16),
            label,
            8.5,
            True,
            color,
            PP_ALIGN.CENTER,
        )


def add_tech_note(slide, text):
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(6.64), Inches(10.95), Inches(0.36))
    fill(note, RGBColor(247, 250, 253))
    note.line.color.rgb = RGBColor(215, 226, 238)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(6.64), Inches(0.07), Inches(0.36))
    fill(bar, BLUE)
    add_text(slide, Inches(0.9), Inches(6.72), Inches(10.55), Inches(0.15), text, 9.4, False, MUTED)


def add_arrow(slide, x1, y1, x2, y2, color=BLUE):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.8)
    return line


def draw_architecture(slide):
    x, y, w, h = 6.55, 2.05, 5.45, 0.72
    layers = [
        ("前端交互层", "HTML / CSS / JS", BLUE),
        ("后端服务层", "FastAPI API", GREEN),
        ("RAG 检索层", "Embedding + Chroma", ORANGE),
        ("本地模型层", "Ollama LLM", RGBColor(105, 91, 180)),
    ]
    for i, (name, sub, color) in enumerate(layers):
        yy = y + i * 0.9
        add_card(slide, x, yy, w, h, name, sub, color)
        if i < len(layers) - 1:
            add_arrow(slide, x + w / 2, yy + h, x + w / 2, yy + 0.9)
    add_card(slide, 6.75, 5.72, 2.05, 0.44, "请求流", "浏览器 → FastAPI → RAG → Ollama", BLUE)
    add_card(slide, 9.2, 5.72, 2.55, 0.44, "数据流", "Markdown → Chunks → Embedding → Vector DB", GREEN)


def draw_pipeline(slide):
    add_picture_card(slide, VISUALS / "knowledge_ingestion.png", 6.18, 2.03, 5.88, 3.28)
    add_label_strip(slide, ["官网资料", "Markdown", "Chunks", "Vector DB"], 6.55, 5.43, 1.12)
    add_card(slide, 6.25, 5.84, 1.55, 0.46, "17 docs", "结构化文档", BLUE)
    add_card(slide, 8.0, 5.84, 1.55, 0.46, "134 chunks", "可检索片段", GREEN)
    add_card(slide, 9.75, 5.84, 2.1, 0.46, "metadata", "标题 / 来源 / 类别", ORANGE)


def draw_rag(slide):
    add_picture_card(slide, VISUALS / "rag_workflow.png", 6.18, 2.03, 5.88, 3.28)
    add_label_strip(slide, ["问题", "向量", "检索", "Prompt", "生成", "来源"], 6.33, 5.43, 0.78)
    add_card(slide, 6.28, 5.84, 2.55, 0.46, "检索增强", "先找资料，再生成答案", GREEN)
    add_card(slide, 9.1, 5.84, 2.55, 0.46, "可追溯", "标题、来源、相似度同步返回", ORANGE)


def draw_fallback(slide):
    add_picture_card(slide, VISUALS / "local_model_fallback.png", 6.18, 2.03, 5.88, 3.28)
    add_card(slide, 6.3, 5.45, 1.42, 0.5, "LLM", "qwen2.5:7b", BLUE)
    add_card(slide, 7.9, 5.45, 1.75, 0.5, "Embedding", "nomic-embed", GREEN)
    add_card(slide, 9.85, 5.45, 1.85, 0.5, "Fallback", "TF-IDF / MOCK", ORANGE)


def draw_frontend(slide):
    x, y = 6.45, 2.05
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.55), Inches(3.85))
    fill(outer, RGBColor(255, 255, 255))
    outer.line.color.rgb = RGBColor(205, 216, 230)
    for i, (name, ww, color) in enumerate([("状态", 1.35, BLUE), ("聊天", 2.15, GREEN), ("流程/来源", 1.7, ORANGE)]):
        xx = x + 0.22 + sum([1.5, 2.3, 1.85][:i])
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xx), Inches(y + 0.38), Inches(ww), Inches(3.05))
        fill(rect, RGBColor(246, 249, 252))
        rect.line.color.rgb = RGBColor(220, 229, 238)
        add_text(slide, Inches(xx + 0.12), Inches(y + 0.55), Inches(ww - 0.22), Inches(0.3), name, 12, True, color)
        for j in range(4):
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(xx + 0.15), Inches(y + 1.0 + j * 0.42), Inches(ww - 0.3), Inches(0.055))
            fill(line, RGBColor(190, 205, 222))
    add_label_strip(slide, ["状态", "问答", "来源", "流程"], 6.7, 5.55, 0.72)


def draw_api(slide):
    rows = [("health", "状态"), ("stats", "统计"), ("rebuild", "索引"), ("chat", "问答")]
    for i, (a, b) in enumerate(rows):
        add_card(slide, 6.6, 2.0 + i * 0.75, 2.3, 0.54, f"/api/{a}", b, BLUE)
    modules = ["main.py", "retriever.py", "prompt_builder.py", "ollama_client.py"]
    for i, mod in enumerate(modules):
        add_card(slide, 9.4, 2.0 + i * 0.75, 2.25, 0.54, mod, "", GREEN)


def draw_validation(slide):
    add_picture_card(slide, VISUALS / "validation_dashboard.png", 6.18, 2.03, 5.88, 3.28)
    add_card(slide, 6.28, 5.45, 1.45, 0.5, "pytest", "8 passed", GREEN)
    add_card(slide, 7.95, 5.45, 1.45, 0.5, "smoke", "PASS", GREEN)
    add_card(slide, 9.62, 5.45, 1.92, 0.5, "acceptance", "PASS", GREEN)


def draw_iteration(slide):
    pairs = [("召回偏差", "混合检索"), ("资料数问题", "运行上下文"), ("重复来源", "文档聚合"), ("索引异常", "清理重建")]
    for i, (p, s) in enumerate(pairs):
        y = 2.0 + i * 0.72
        add_card(slide, 6.45, y, 1.75, 0.5, p, "", ORANGE)
        add_arrow(slide, 8.25, y + 0.25, 8.75, y + 0.25)
        add_card(slide, 8.9, y, 2.35, 0.5, s, "", GREEN)


def draw_results(slide):
    metrics = [("17", "知识库文档"), ("134", "chunks"), ("4", "核心接口"), ("8", "单元测试")]
    for i, (num, label) in enumerate(metrics):
        x = 6.4 + (i % 2) * 2.6
        y = 2.2 + (i // 2) * 1.45
        add_card(slide, x, y, 2.25, 1.05, num, label, [BLUE, GREEN, ORANGE, RGBColor(105, 91, 180)][i])
        # Make metric number bigger by overlay.
        add_text(slide, Inches(x + 0.22), Inches(y + 0.12), Inches(0.82), Inches(0.4), num, 22, True, [BLUE, GREEN, ORANGE, RGBColor(105, 91, 180)][i])


def draw_generic_cards(slide, bullets):
    for i, bullet in enumerate(bullets[:4]):
        x = 6.45 + (i % 2) * 2.75
        y = 2.15 + (i // 2) * 1.55
        title = bullet.split("：", 1)[0] if "：" in bullet else bullet[:10]
        body = bullet.split("：", 1)[1] if "：" in bullet else ""
        add_card(slide, x, y, 2.35, 1.08, title, body, [BLUE, GREEN, ORANGE, RGBColor(105, 91, 180)][i])


def make_body_slide(prs, idx, section, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[len(prs.slide_layouts) - 1])
    clear_slide(slide)
    add_header(slide, idx, section, title)
    bullet_size = 15.0 if len(bullets) >= 6 else 15.6 if len(bullets) >= 5 else 16.6
    add_bullet_list(slide, bullets, x=0.75, y=2.16, w=5.15, h=3.95, size=bullet_size)
    if idx == 5:
        draw_architecture(slide)
    elif idx == 6:
        draw_pipeline(slide)
    elif idx == 7:
        draw_rag(slide)
    elif idx == 8:
        draw_fallback(slide)
    elif idx == 9:
        draw_frontend(slide)
    elif idx == 10:
        draw_api(slide)
    elif idx == 11:
        draw_validation(slide)
    elif idx == 12:
        draw_iteration(slide)
    elif idx == 13:
        draw_results(slide)
    else:
        draw_generic_cards(slide, bullets)
    if idx in TECH_NOTES:
        add_tech_note(slide, TECH_NOTES[idx])


def build() -> None:
    if CURRENT.exists() and not CLEAN_BACKUP.exists():
        shutil.copy2(CURRENT, CLEAN_BACKUP)
    if CURRENT.exists() and not OLD_BUSY.exists():
        shutil.copy2(CURRENT, OLD_BUSY)
    prs = Presentation(str(TEMPLATE_BACKUP))
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    delete_all_slides(prs)
    blank = prs.slide_layouts[len(prs.slide_layouts) - 1]

    # Cover.
    slide = prs.slides.add_slide(blank)
    clear_slide(slide)
    add_logo(slide, 0.35, 0.2, 1.75)
    add_anniversary(slide, 11.25, 0.1, 0.95)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.85), Inches(13.333), Inches(2.0))
    fill(banner, BLUE)
    add_text(slide, Inches(1.1), Inches(2.22), Inches(11.2), Inches(0.9), SLIDES[0][0], 31, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)
    add_text(slide, Inches(1.1), Inches(3.18), Inches(11.2), Inches(0.28), SLIDES[0][1], 12, False, RGBColor(255, 255, 255), PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(5.25), Inches(13.333), Inches(0.42), "天津大学国际工程师学院", 18, True, BLUE_DARK, PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(6.05), Inches(13.333), Inches(0.3), SLIDES[0][2], 12, True, BLUE_DARK, PP_ALIGN.CENTER)

    # Outline.
    slide = prs.slides.add_slide(blank)
    clear_slide(slide)
    blue_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.35), Inches(7.5))
    fill(blue_block, BLUE)
    add_text(slide, Inches(0.92), Inches(3.28), Inches(1.9), Inches(0.46), "OUTLINE", 25, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)
    for i, item in enumerate(SLIDES[1][2], 1):
        y = 1.55 + (i - 1) * 1.45
        num = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.35), Inches(y), Inches(0.45), Inches(0.45))
        fill(num, BLUE)
        add_text(slide, Inches(4.35), Inches(y + 0.06), Inches(0.45), Inches(0.24), str(i), 13, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)
        add_text(slide, Inches(5.15), Inches(y + 0.02), Inches(5.8), Inches(0.35), item, 18, True, BLUE_DARK)

    for idx, item in enumerate(SLIDES[2:], 3):
        section, title, bullets = item
        if idx == 14:
            slide = prs.slides.add_slide(blank)
            clear_slide(slide)
            add_logo(slide, 0.35, 0.2, 1.55)
            add_text(slide, Inches(0), Inches(2.45), Inches(13.333), Inches(0.62), "THANK YOU !", 34, True, BLUE_DARK, PP_ALIGN.CENTER)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(3.25), Inches(7.3), Inches(0.035))
            fill(line, BLUE)
            add_text(slide, Inches(2.4), Inches(3.58), Inches(8.5), Inches(0.42), "总结与展望：从课堂 Demo 到企业知识服务原型", 18, True, TEXT, PP_ALIGN.CENTER)
            add_bullet_list(slide, bullets, x=2.55, y=4.18, w=8.1, h=1.65, size=13.5)
            add_label_strip(slide, ["资料入库", "向量检索", "本地生成", "来源追溯", "持续迭代"], 3.28, 6.02, 1.08)
            continue
        make_body_slide(prs, idx, section, title, bullets)

    prs.save(DONE)
    shutil.copy2(DONE, CURRENT)
    print(DONE)
    print(CURRENT)
    print(OLD_BUSY)


if __name__ == "__main__":
    build()
